from __future__ import annotations

import csv
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile

from docx import Document
from openpyxl import Workbook
from PIL import Image, ImageDraw
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas


def create_text_pdf(path: Path) -> None:
    pdf = canvas.Canvas(str(path), pagesize=letter)
    pdf.setTitle("Smoke Text PDF")
    pdf.drawString(72, 720, "MarkItDown smoke test PDF")
    pdf.drawString(72, 700, "This PDF contains an actual text layer.")
    pdf.save()


def create_scanned_pdf(path: Path, temp_image_path: Path) -> None:
    image = Image.new("RGB", (1000, 700), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle((40, 40, 960, 660), outline="black", width=4)
    draw.text((80, 120), "SCANNED PAGE", fill="black")
    draw.text((80, 180), "This text exists only inside the image.", fill="black")
    image.save(temp_image_path)

    pdf = canvas.Canvas(str(path), pagesize=letter)
    pdf.drawImage(ImageReader(str(temp_image_path)), 36, 180, width=540, preserveAspectRatio=True, mask="auto")
    pdf.save()


def create_docx(path: Path) -> None:
    document = Document()
    document.add_heading("MarkItDown Smoke DOCX", level=1)
    document.add_paragraph("This DOCX verifies local-only conversion.")
    document.add_paragraph("Bullet one", style="List Bullet")
    document.add_paragraph("Bullet two", style="List Bullet")
    document.save(path)


def create_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "MarkItDown Smoke PPTX"
    slide.placeholders[1].text = "Local conversion only\nNo remote fetches"
    presentation.save(path)


def create_xlsx(path: Path) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Smoke"
    sheet.append(["Format", "Expectation"])
    sheet.append(["xlsx", "PASS"])
    sheet.append(["pdf-scan", "NEEDS_OCR"])
    workbook.save(path)


def create_html(path: Path) -> None:
    path.write_text(
        "<html><body><h1>MarkItDown Smoke HTML</h1><p>Simple local HTML sample.</p></body></html>",
        encoding="utf-8",
    )


def create_csv(path: Path) -> None:
    with path.open("w", newline="", encoding="utf-8") as handle:
        writer = csv.writer(handle)
        writer.writerow(["format", "status"])
        writer.writerow(["csv", "PASS"])


def create_zip(path: Path) -> None:
    with ZipFile(path, "w", compression=ZIP_DEFLATED) as archive:
        archive.writestr("inside/readme.txt", "This ZIP file is part of the MarkItDown smoke suite.")


def generate_samples(output_dir: Path) -> dict[str, Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    samples = {
        "pdf_text": output_dir / "sample-text.pdf",
        "pdf_scan": output_dir / "sample-scan.pdf",
        "docx": output_dir / "sample.docx",
        "pptx": output_dir / "sample.pptx",
        "xlsx": output_dir / "sample.xlsx",
        "html": output_dir / "sample.html",
        "csv": output_dir / "sample.csv",
        "zip": output_dir / "sample.zip",
    }

    create_text_pdf(samples["pdf_text"])
    temp_image_path = output_dir / "sample-scan-source.png"
    create_scanned_pdf(samples["pdf_scan"], temp_image_path)
    create_docx(samples["docx"])
    create_pptx(samples["pptx"])
    create_xlsx(samples["xlsx"])
    create_html(samples["html"])
    create_csv(samples["csv"])
    create_zip(samples["zip"])
    temp_image_path.unlink(missing_ok=True)
    return samples
